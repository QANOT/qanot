"""Document generation tools — Word, Excel, PDF, and PowerPoint."""

from __future__ import annotations

import json
import logging
import re
from pathlib import Path

from qanot.agent import ToolRegistry

logger = logging.getLogger(__name__)


def _parse_page_range(pages_str: str, total: int) -> list[int]:
    """Parse '1-3' or '1,3,5' into zero-based page indices."""
    indices: list[int] = []
    for part in pages_str.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            try:
                s, e = int(start) - 1, int(end) - 1
                indices.extend(range(max(0, s), min(total, e + 1)))
            except ValueError:
                pass
        else:
            try:
                indices.append(int(part) - 1)
            except ValueError:
                pass
    return indices


def register_document_tools(registry: ToolRegistry, workspace_dir: str) -> None:
    """Register document creation tools."""

    # ── create_docx ──
    async def create_docx(params: dict) -> str:
        """Create a Word document."""
        try:
            from docx import Document
            from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
        except ImportError:
            return json.dumps({"error": "python-docx kutubxonasi o'rnatilmagan. pip install python-docx"})

        filename = params.get("filename", "document.docx")
        if not filename.endswith(".docx"):
            filename += ".docx"

        title = params.get("title", "")
        content = params.get("content", "")
        rows = params.get("rows")

        doc = Document()

        # Title
        if title:
            p = doc.add_heading(title, level=0)
            p.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        # Content — parse simple markdown-like format
        if content:
            for line in content.split("\n"):
                line = line.strip()
                if not line:
                    doc.add_paragraph("")
                elif line.startswith("# "):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith("## "):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith("### "):
                    doc.add_heading(line[4:], level=3)
                elif line.startswith("- ") or line.startswith("* "):
                    doc.add_paragraph(line[2:], style="List Bullet")
                elif re.match(r"^\d+\.\s", line):
                    doc.add_paragraph(re.sub(r"^\d+\.\s*", "", line), style="List Number")
                else:
                    p = doc.add_paragraph()
                    _add_formatted_text(p, line)

        # Table
        if rows and isinstance(rows, list) and len(rows) > 0:
            table = doc.add_table(rows=len(rows), cols=len(rows[0]))
            table.style = "Table Grid"
            for i, row in enumerate(rows):
                for j, cell_text in enumerate(row):
                    table.rows[i].cells[j].text = str(cell_text)
                    if i == 0:
                        for paragraph in table.rows[i].cells[j].paragraphs:
                            for run in paragraph.runs:
                                run.bold = True

        # Save
        filepath = Path(workspace_dir) / filename
        doc.save(str(filepath))

        return json.dumps({
            "status": "ok",
            "file": str(filepath),
            "filename": filename,
            "message": f"{filename} yaratildi",
        })

    def _add_formatted_text(paragraph, text: str) -> None:
        """Parse simple bold (**text**) formatting."""
        parts = re.split(r"(\*\*.*?\*\*)", text)
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                run = paragraph.add_run(part[2:-2])
                run.bold = True
            else:
                paragraph.add_run(part)

    # ── create_xlsx ──
    async def create_xlsx(params: dict) -> str:
        """Create an Excel spreadsheet."""
        try:
            from openpyxl import Workbook
        except ImportError:
            return json.dumps({"error": "openpyxl kutubxonasi o'rnatilmagan. pip install openpyxl"})

        filename = params.get("filename", "spreadsheet.xlsx")
        if not filename.endswith(".xlsx"):
            filename += ".xlsx"

        sheets = params.get("sheets")
        headers = params.get("headers")
        rows = params.get("rows")
        title = params.get("title", "")

        wb = Workbook()

        if sheets and isinstance(sheets, list):
            for idx, sheet_data in enumerate(sheets):
                if idx == 0:
                    ws = wb.active
                else:
                    ws = wb.create_sheet()
                ws.title = sheet_data.get("name", f"Sheet{idx + 1}")
                _fill_sheet(
                    ws,
                    sheet_data.get("headers", []),
                    sheet_data.get("rows", []),
                    sheet_data.get("title", ""),
                )
        else:
            ws = wb.active
            ws.title = title or "Sheet1"
            _fill_sheet(ws, headers or [], rows or [], title)

        filepath = Path(workspace_dir) / filename
        wb.save(str(filepath))

        return json.dumps({
            "status": "ok",
            "file": str(filepath),
            "filename": filename,
            "message": f"{filename} yaratildi",
        })

    def _fill_sheet(ws, headers: list, rows: list, title: str = "") -> None:
        """Fill a worksheet with headers and rows."""
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

        start_row = 1

        # Title row
        if title:
            ws.cell(row=1, column=1, value=title)
            ws.cell(row=1, column=1).font = Font(bold=True, size=14)
            start_row = 3

        # Header row
        if headers:
            header_fill = PatternFill(start_color="2563EB", end_color="2563EB", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF", size=11)
            thin_border = Border(
                left=Side(style="thin"),
                right=Side(style="thin"),
                top=Side(style="thin"),
                bottom=Side(style="thin"),
            )

            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=start_row, column=col_idx, value=str(header))
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")
                cell.border = thin_border
                ws.column_dimensions[cell.column_letter].width = max(len(str(header)) + 4, 12)
            start_row += 1

        # Data rows
        if rows:
            thin_border = Border(
                left=Side(style="thin"),
                right=Side(style="thin"),
                top=Side(style="thin"),
                bottom=Side(style="thin"),
            )
            for row_idx, row_data in enumerate(rows, start_row):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.border = thin_border
                    # Auto-detect numbers
                    if isinstance(value, str):
                        try:
                            cell.value = float(value.replace(",", "").replace(" ", ""))
                            cell.number_format = "#,##0"
                        except (ValueError, AttributeError):
                            pass

    # Register tools
    registry.register(
        name="create_docx",
        description=(
            "Word (.docx) hujjat yaratish. Shartnoma, hisobot, taklifnoma va boshqa hujjatlar uchun. "
            "Markdown formatda content yozing: # sarlavha, ## kichik sarlavha, **qalin**, - ro'yxat. "
            "Jadval uchun rows parametrini ishlating."
        ),
        parameters={
            "type": "object",
            "required": ["filename"],
            "properties": {
                "filename": {"type": "string", "description": "Fayl nomi (masalan: shartnoma.docx)"},
                "title": {"type": "string", "description": "Hujjat sarlavhasi"},
                "content": {"type": "string", "description": "Hujjat matni (Markdown format)"},
                "rows": {
                    "type": "array",
                    "description": "Jadval ma'lumotlari. Birinchi qator — sarlavha. [[\"Nomi\",\"Narxi\"],[\"Tovar1\",\"5000\"]]",
                    "items": {"type": "array", "items": {"type": "string"}},
                },
            },
        },
        handler=create_docx,
    )

    registry.register(
        name="create_xlsx",
        description=(
            "Excel (.xlsx) jadval yaratish. Hisobot, ro'yxat, statistika uchun. "
            "headers va rows bering — chiroyli formatlanadi. "
            "Ko'p sahifali Excel uchun sheets parametrini ishlating."
        ),
        parameters={
            "type": "object",
            "required": ["filename"],
            "properties": {
                "filename": {"type": "string", "description": "Fayl nomi (masalan: hisobot.xlsx)"},
                "title": {"type": "string", "description": "Sarlavha (birinchi qator)"},
                "headers": {
                    "type": "array",
                    "description": "Ustun sarlavhalari: [\"Sana\", \"Summa\", \"Mijoz\"]",
                    "items": {"type": "string"},
                },
                "rows": {
                    "type": "array",
                    "description": "Ma'lumot qatorlari: [[\"2026-03-16\", 5000000, \"Sardor\"]]",
                    "items": {"type": "array"},
                },
                "sheets": {
                    "type": "array",
                    "description": "Ko'p sahifa: [{name, title, headers, rows}]",
                    "items": {"type": "object"},
                },
            },
        },
        handler=create_xlsx,
    )

    # ── read_docx ──
    async def read_docx(params: dict) -> str:
        """Read content from a Word document."""
        try:
            from docx import Document
        except ImportError:
            return json.dumps({"error": "python-docx kutubxonasi o'rnatilmagan"})

        filepath = params.get("file", "")
        if not filepath:
            return json.dumps({"error": "file parametri kerak"})

        path = Path(filepath) if Path(filepath).is_absolute() else Path(workspace_dir) / filepath
        if not path.exists():
            return json.dumps({"error": f"Fayl topilmadi: {filepath}"})

        doc = Document(str(path))
        content = []
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading ", "").replace("Heading", "1")
                try:
                    level = int(level)
                except ValueError:
                    level = 1
                content.append(f"{'#' * level} {para.text}")
            elif para.style.name == "List Bullet":
                content.append(f"- {para.text}")
            elif para.style.name == "List Number":
                content.append(f"1. {para.text}")
            else:
                content.append(para.text)

        # Read tables
        tables_text = []
        for i, table in enumerate(doc.tables):
            table_rows = []
            for row in table.rows:
                table_rows.append([cell.text for cell in row.cells])
            tables_text.append({"table_index": i, "rows": table_rows})

        return json.dumps({
            "content": "\n".join(content),
            "tables": tables_text,
            "paragraphs": len(doc.paragraphs),
        }, ensure_ascii=False)

    # ── edit_docx ──
    async def edit_docx(params: dict) -> str:
        """Edit an existing Word document — append text, replace text, or add table."""
        try:
            from docx import Document
            from docx.shared import Pt  # noqa: F841
        except ImportError:
            return json.dumps({"error": "python-docx kutubxonasi o'rnatilmagan"})

        filepath = params.get("file", "")
        if not filepath:
            return json.dumps({"error": "file parametri kerak"})

        path = Path(filepath) if Path(filepath).is_absolute() else Path(workspace_dir) / filepath
        if not path.exists():
            return json.dumps({"error": f"Fayl topilmadi: {filepath}"})

        doc = Document(str(path))
        action = params.get("action", "append")  # append, replace, add_table

        if action == "append":
            content = params.get("content", "")
            for line in content.split("\n"):
                line = line.strip()
                if not line:
                    doc.add_paragraph("")
                elif line.startswith("# "):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith("## "):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith("- "):
                    doc.add_paragraph(line[2:], style="List Bullet")
                else:
                    p = doc.add_paragraph()
                    _add_formatted_text(p, line)

        elif action == "replace":
            old_text = params.get("old_text", "")
            new_text = params.get("new_text", "")
            if not old_text:
                return json.dumps({"error": "old_text kerak"})
            count = 0
            for para in doc.paragraphs:
                if old_text in para.text:
                    for run in para.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            count += 1
            if count == 0:
                return json.dumps({"error": f"'{old_text}' topilmadi"})

        elif action == "add_table":
            rows = params.get("rows", [])
            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                table.style = "Table Grid"
                for i, row in enumerate(rows):
                    for j, val in enumerate(row):
                        table.rows[i].cells[j].text = str(val)

        doc.save(str(path))
        return json.dumps({"status": "ok", "message": f"{path.name} yangilandi", "action": action})

    # ── read_xlsx ──
    async def read_xlsx(params: dict) -> str:
        """Read data from Excel spreadsheet."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            return json.dumps({"error": "openpyxl kutubxonasi o'rnatilmagan"})

        filepath = params.get("file", "")
        if not filepath:
            return json.dumps({"error": "file parametri kerak"})

        path = Path(filepath) if Path(filepath).is_absolute() else Path(workspace_dir) / filepath
        if not path.exists():
            return json.dumps({"error": f"Fayl topilmadi: {filepath}"})

        wb = load_workbook(str(path), read_only=True, data_only=True)
        sheet_name = params.get("sheet")  # None = active sheet
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active

        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append([str(cell) if cell is not None else "" for cell in row])

        sheet_names = wb.sheetnames
        sheet_title = ws.title
        wb.close()

        return json.dumps({
            "sheet": sheet_title,
            "sheets": sheet_names,
            "rows": rows[:500],
            "total_rows": len(rows),
            "truncated": len(rows) > 500,
        }, ensure_ascii=False)

    # ── edit_xlsx ──
    async def edit_xlsx(params: dict) -> str:
        """Edit an existing Excel spreadsheet — append rows, update cells, add sheet."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            return json.dumps({"error": "openpyxl kutubxonasi o'rnatilmagan"})

        filepath = params.get("file", "")
        if not filepath:
            return json.dumps({"error": "file parametri kerak"})

        path = Path(filepath) if Path(filepath).is_absolute() else Path(workspace_dir) / filepath
        if not path.exists():
            return json.dumps({"error": f"Fayl topilmadi: {filepath}"})

        wb = load_workbook(str(path))
        sheet_name = params.get("sheet")
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active

        action = params.get("action", "append_rows")

        if action == "append_rows":
            rows = params.get("rows", [])
            for row_data in rows:
                ws.append(row_data)

        elif action == "update_cell":
            cell = params.get("cell", "")  # e.g. "A1", "B5"
            value = params.get("value", "")
            if cell:
                ws[cell] = value

        elif action == "add_sheet":
            new_name = params.get("new_sheet_name", "Sheet")
            ws_new = wb.create_sheet(title=new_name)
            headers = params.get("headers", [])
            if headers:
                ws_new.append(headers)

        wb.save(str(path))
        wb.close()

        return json.dumps({"status": "ok", "message": f"{path.name} yangilandi", "action": action})

    # Register read/edit tools
    registry.register(
        name="read_docx",
        description="Word (.docx) hujjatni o'qish. Matn, sarlavhalar, jadvallarni qaytaradi.",
        parameters={
            "type": "object",
            "required": ["file"],
            "properties": {
                "file": {"type": "string", "description": "Fayl nomi yoki to'liq path"},
            },
        },
        handler=read_docx,
    )

    registry.register(
        name="edit_docx",
        description=(
            "Word (.docx) hujjatni tahrirlash. Matn qo'shish (append), "
            "almashtirish (replace), jadval qo'shish (add_table)."
        ),
        parameters={
            "type": "object",
            "required": ["file"],
            "properties": {
                "file": {"type": "string", "description": "Fayl nomi yoki to'liq path"},
                "action": {
                    "type": "string",
                    "enum": ["append", "replace", "add_table"],
                    "description": "Harakat turi: append (matn qo'shish), replace (almashtirish), add_table (jadval)",
                },
                "content": {"type": "string", "description": "Qo'shiladigan matn (append uchun, Markdown format)"},
                "old_text": {"type": "string", "description": "Almashtirilishi kerak bo'lgan matn (replace uchun)"},
                "new_text": {"type": "string", "description": "Yangi matn (replace uchun)"},
                "rows": {
                    "type": "array",
                    "description": "Jadval qatorlari (add_table uchun)",
                    "items": {"type": "array", "items": {"type": "string"}},
                },
            },
        },
        handler=edit_docx,
    )

    registry.register(
        name="read_xlsx",
        description="Excel (.xlsx) fayldan ma'lumot o'qish. Barcha qatorlar va ustunlarni qaytaradi.",
        parameters={
            "type": "object",
            "required": ["file"],
            "properties": {
                "file": {"type": "string", "description": "Fayl nomi yoki to'liq path"},
                "sheet": {"type": "string", "description": "Sahifa nomi (ixtiyoriy, standart — faol sahifa)"},
            },
        },
        handler=read_xlsx,
    )

    registry.register(
        name="edit_xlsx",
        description=(
            "Excel (.xlsx) faylni tahrirlash. Qator qo'shish (append_rows), "
            "katak yangilash (update_cell), yangi sahifa (add_sheet)."
        ),
        parameters={
            "type": "object",
            "required": ["file"],
            "properties": {
                "file": {"type": "string", "description": "Fayl nomi yoki to'liq path"},
                "sheet": {"type": "string", "description": "Sahifa nomi (ixtiyoriy)"},
                "action": {
                    "type": "string",
                    "enum": ["append_rows", "update_cell", "add_sheet"],
                    "description": "Harakat turi",
                },
                "rows": {
                    "type": "array",
                    "description": "Qo'shiladigan qatorlar (append_rows uchun)",
                    "items": {"type": "array"},
                },
                "cell": {"type": "string", "description": "Katak manzili, masalan A1, B5 (update_cell uchun)"},
                "value": {"type": "string", "description": "Yangi qiymat (update_cell uchun)"},
                "new_sheet_name": {"type": "string", "description": "Yangi sahifa nomi (add_sheet uchun)"},
                "headers": {
                    "type": "array",
                    "description": "Ustun sarlavhalari (add_sheet uchun)",
                    "items": {"type": "string"},
                },
            },
        },
        handler=edit_xlsx,
    )

    # ── create_pdf ──
    async def create_pdf(params: dict) -> str:
        """Create a PDF document."""
        try:
            from fpdf import FPDF
        except ImportError:
            return json.dumps({"error": "fpdf2 kutubxonasi o'rnatilmagan. pip install fpdf2"})

        filename = params.get("filename", "document.pdf")
        if not filename.endswith(".pdf"):
            filename += ".pdf"

        title = params.get("title", "")
        content = params.get("content", "")
        rows = params.get("rows")

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()

        # Use built-in font (supports basic latin)
        pdf.set_font("Helvetica", size=12)

        # Title
        if title:
            pdf.set_font("Helvetica", "B", 18)
            pdf.cell(0, 15, title, new_x="LMARGIN", new_y="NEXT", align="C")
            pdf.ln(5)
            pdf.set_font("Helvetica", size=12)

        # Content — parse line by line
        if content:
            for line in content.split("\n"):
                line = line.strip()
                if not line:
                    pdf.ln(5)
                elif line.startswith("# "):
                    pdf.set_font("Helvetica", "B", 16)
                    pdf.cell(0, 10, line[2:], new_x="LMARGIN", new_y="NEXT")
                    pdf.set_font("Helvetica", size=12)
                elif line.startswith("## "):
                    pdf.set_font("Helvetica", "B", 14)
                    pdf.cell(0, 9, line[3:], new_x="LMARGIN", new_y="NEXT")
                    pdf.set_font("Helvetica", size=12)
                elif line.startswith("- ") or line.startswith("* "):
                    pdf.cell(10)
                    pdf.cell(0, 7, f"\u2022 {line[2:]}", new_x="LMARGIN", new_y="NEXT")
                else:
                    pdf.multi_cell(0, 7, line)

        # Table
        if rows and isinstance(rows, list) and len(rows) > 0:
            pdf.ln(5)
            col_count = len(rows[0])
            col_width = (pdf.w - 20) / col_count

            # Header row
            pdf.set_font("Helvetica", "B", 11)
            pdf.set_fill_color(37, 99, 235)
            pdf.set_text_color(255, 255, 255)
            for cell in rows[0]:
                pdf.cell(col_width, 8, str(cell), border=1, fill=True, align="C")
            pdf.ln()

            # Data rows
            pdf.set_font("Helvetica", size=10)
            pdf.set_text_color(0, 0, 0)
            for row in rows[1:]:
                for cell in row:
                    pdf.cell(col_width, 7, str(cell), border=1)
                pdf.ln()

        filepath = Path(workspace_dir) / filename
        pdf.output(str(filepath))

        return json.dumps({
            "status": "ok",
            "file": str(filepath),
            "filename": filename,
            "message": f"{filename} yaratildi",
        })

    # ── read_pdf ──
    async def read_pdf(params: dict) -> str:
        """Read content from a PDF file."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return json.dumps({"error": "PyMuPDF kutubxonasi o'rnatilmagan. pip install PyMuPDF"})

        filepath = params.get("file", "")
        if not filepath:
            return json.dumps({"error": "file parametri kerak"})

        path = Path(filepath) if Path(filepath).is_absolute() else Path(workspace_dir) / filepath
        if not path.exists():
            return json.dumps({"error": f"Fayl topilmadi: {filepath}"})

        doc = fitz.open(str(path))
        total_pages = len(doc)

        # Parse page range
        pages = params.get("pages")
        page_indices = list(range(total_pages))
        if pages:
            page_indices = _parse_page_range(pages, total_pages)

        content = []
        for i in page_indices:
            if 0 <= i < total_pages:
                page = doc[i]
                text = page.get_text()
                content.append(f"--- Sahifa {i + 1} ---\n{text}")

        doc.close()

        full_text = "\n\n".join(content)
        if len(full_text) > 50000:
            full_text = full_text[:50000] + "\n\n[Truncated \u2014 juda katta fayl]"

        return json.dumps({
            "content": full_text,
            "total_pages": total_pages,
            "pages_read": len(page_indices),
        }, ensure_ascii=False)

    # ── create_pptx ──
    async def create_pptx(params: dict) -> str:
        """Create a PowerPoint presentation."""
        try:
            from pptx import Presentation
        except ImportError:
            return json.dumps({"error": "python-pptx kutubxonasi o'rnatilmagan. pip install python-pptx"})

        filename = params.get("filename", "presentation.pptx")
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        title = params.get("title", "Prezentatsiya")
        slides = params.get("slides", [])

        prs = Presentation()

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = params.get("subtitle", "")

        # Content slides
        for slide_data in slides:
            if isinstance(slide_data, str):
                parts = slide_data.split("\n")
                slide_data = {"title": parts[0], "content": "\n".join(parts[1:])}

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            slide.shapes.title.text = slide_data.get("title", "")

            content = slide_data.get("content", "")
            if content and slide.placeholders[1]:
                tf = slide.placeholders[1].text_frame
                tf.text = ""
                first_line = True
                for line in content.split("\n"):
                    line = line.strip()
                    if not line:
                        continue
                    if first_line:
                        tf.text = line.lstrip("- *\u2022")
                        first_line = False
                    else:
                        p = tf.add_paragraph()
                        p.text = line.lstrip("- *\u2022")
                        if line.startswith(("- ", "* ", "\u2022 ")):
                            p.level = 1

        filepath = Path(workspace_dir) / filename
        prs.save(str(filepath))

        return json.dumps({
            "status": "ok",
            "file": str(filepath),
            "filename": filename,
            "slides_count": len(prs.slides),
            "message": f"{filename} yaratildi ({len(prs.slides)} slayd)",
        })

    # ── read_pptx ──
    async def read_pptx(params: dict) -> str:
        """Read content from a PowerPoint presentation."""
        try:
            from pptx import Presentation
        except ImportError:
            return json.dumps({"error": "python-pptx kutubxonasi o'rnatilmagan. pip install python-pptx"})

        filepath = params.get("file", "")
        if not filepath:
            return json.dumps({"error": "file parametri kerak"})

        path = Path(filepath) if Path(filepath).is_absolute() else Path(workspace_dir) / filepath
        if not path.exists():
            return json.dumps({"error": f"Fayl topilmadi: {filepath}"})

        prs = Presentation(str(path))
        slides_data = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text)
            slides_data.append({
                "slide": i + 1,
                "content": "\n".join(slide_text),
            })

        return json.dumps({
            "total_slides": len(slides_data),
            "slides": slides_data,
        }, ensure_ascii=False)

    # Register PDF tools
    registry.register(
        name="create_pdf",
        description=(
            "PDF hujjat yaratish. Shartnoma, hisobot, faktura uchun. "
            "Markdown format qo'llanadi: # sarlavha, ## kichik sarlavha, - ro'yxat. "
            "Jadval uchun rows parametrini ishlating."
        ),
        parameters={
            "type": "object",
            "required": ["filename"],
            "properties": {
                "filename": {"type": "string", "description": "Fayl nomi (masalan: hisobot.pdf)"},
                "title": {"type": "string", "description": "Hujjat sarlavhasi"},
                "content": {"type": "string", "description": "Hujjat matni (Markdown format)"},
                "rows": {
                    "type": "array",
                    "description": "Jadval ma'lumotlari. Birinchi qator — sarlavha.",
                    "items": {"type": "array", "items": {"type": "string"}},
                },
            },
        },
        handler=create_pdf,
    )

    registry.register(
        name="read_pdf",
        description="PDF faylni o'qish. Barcha yoki tanlangan sahifalar matnini qaytaradi.",
        parameters={
            "type": "object",
            "required": ["file"],
            "properties": {
                "file": {"type": "string", "description": "Fayl nomi yoki to'liq path"},
                "pages": {
                    "type": "string",
                    "description": "Sahifa diapazoni: '1-3' yoki '1,3,5' (ixtiyoriy, standart — hammasi)",
                },
            },
        },
        handler=read_pdf,
    )

    # Register PowerPoint tools
    registry.register(
        name="create_pptx",
        description=(
            "PowerPoint (.pptx) prezentatsiya yaratish. Sarlavha va slaydlar bering. "
            "Har bir slayd title va content dan iborat."
        ),
        parameters={
            "type": "object",
            "required": ["filename"],
            "properties": {
                "filename": {"type": "string", "description": "Fayl nomi (masalan: taqdimot.pptx)"},
                "title": {"type": "string", "description": "Prezentatsiya sarlavhasi"},
                "subtitle": {"type": "string", "description": "Sarlavha ostidagi matn"},
                "slides": {
                    "type": "array",
                    "description": "Slaydlar ro'yxati: [{title, content}]",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "content": {"type": "string"},
                        },
                    },
                },
            },
        },
        handler=create_pptx,
    )

    registry.register(
        name="read_pptx",
        description="PowerPoint (.pptx) prezentatsiyani o'qish. Barcha slaydlar matnini qaytaradi.",
        parameters={
            "type": "object",
            "required": ["file"],
            "properties": {
                "file": {"type": "string", "description": "Fayl nomi yoki to'liq path"},
            },
        },
        handler=read_pptx,
    )

    logger.info(
        "Document tools registered: create_docx, create_xlsx, read_docx, edit_docx, "
        "read_xlsx, edit_xlsx, create_pdf, read_pdf, create_pptx, read_pptx"
    )
