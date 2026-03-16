"""PowerPoint (.pptx) presentation tools — create, read, edit."""

from __future__ import annotations

import json
import logging
from pathlib import Path

from qanot.registry import ToolRegistry

logger = logging.getLogger(__name__)


def register_pptx_tools(registry: ToolRegistry, workspace_dir: str) -> None:
    """Register PowerPoint presentation tools."""

    # ── create_pptx ──
    async def create_pptx(params: dict) -> str:
        """Create a PowerPoint presentation."""
        try:
            from pptx import Presentation
        except ImportError:
            return json.dumps({"error": "python-pptx kutubxonasi o'rnatilmagan. pip install python-pptx"})

        filename = params.get("filename", "presentation.pptx")
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        title = params.get("title", "Prezentatsiya")
        slides = params.get("slides", [])

        prs = Presentation()

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = params.get("subtitle", "")

        # Content slides
        for slide_data in slides:
            if isinstance(slide_data, str):
                parts = slide_data.split("\n")
                slide_data = {"title": parts[0], "content": "\n".join(parts[1:])}

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            slide.shapes.title.text = slide_data.get("title", "")

            content = slide_data.get("content", "")
            if content and slide.placeholders[1]:
                tf = slide.placeholders[1].text_frame
                tf.text = ""
                first_line = True
                for line in content.split("\n"):
                    line = line.strip()
                    if not line:
                        continue
                    if first_line:
                        tf.text = line.lstrip("- *\u2022")
                        first_line = False
                    else:
                        p = tf.add_paragraph()
                        p.text = line.lstrip("- *\u2022")
                        if line.startswith(("- ", "* ", "\u2022 ")):
                            p.level = 1

        filepath = Path(workspace_dir) / filename
        prs.save(str(filepath))

        return json.dumps({
            "status": "ok",
            "file": str(filepath),
            "filename": filename,
            "slides_count": len(prs.slides),
            "message": f"{filename} yaratildi ({len(prs.slides)} slayd)",
        })

    # ── read_pptx ──
    async def read_pptx(params: dict) -> str:
        """Read content from a PowerPoint presentation."""
        try:
            from pptx import Presentation
        except ImportError:
            return json.dumps({"error": "python-pptx kutubxonasi o'rnatilmagan. pip install python-pptx"})

        filepath = params.get("file", "")
        if not filepath:
            return json.dumps({"error": "file parametri kerak"})

        path = Path(filepath) if Path(filepath).is_absolute() else Path(workspace_dir) / filepath
        if not path.exists():
            return json.dumps({"error": f"Fayl topilmadi: {filepath}"})

        prs = Presentation(str(path))
        slides_data = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text)
            slides_data.append({
                "slide": i + 1,
                "content": "\n".join(slide_text),
            })

        return json.dumps({
            "total_slides": len(slides_data),
            "slides": slides_data,
        }, ensure_ascii=False)

    # ── edit_pptx ──
    async def edit_pptx(params: dict) -> str:
        """Edit existing PowerPoint — add slides, replace text, delete slides."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt  # noqa: F401
        except ImportError:
            return json.dumps({"error": "python-pptx kutubxonasi o'rnatilmagan"})

        filepath = params.get("file", "")
        path = Path(filepath) if Path(filepath).is_absolute() else Path(workspace_dir) / filepath
        if not path.exists():
            return json.dumps({"error": f"Fayl topilmadi: {filepath}"})

        prs = Presentation(str(path))
        action = params.get("action", "add_slide")

        if action == "add_slide":
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = params.get("title", "")
            content = params.get("content", "")
            if content and slide.placeholders[1]:
                tf = slide.placeholders[1].text_frame
                tf.text = ""
                for i, line in enumerate(content.split("\n")):
                    line = line.strip()
                    if not line:
                        continue
                    if i == 0:
                        tf.text = line.lstrip("- *")
                    else:
                        p = tf.add_paragraph()
                        p.text = line.lstrip("- *")

        elif action == "replace":
            old_text = params.get("old_text", "")
            new_text = params.get("new_text", "")
            if not old_text:
                return json.dumps({"error": "old_text kerak"})
            count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)
                                    count += 1
            if count == 0:
                return json.dumps({"error": f"'{old_text}' topilmadi"})

        elif action == "delete_slide":
            slide_index = params.get("slide_index", -1)
            if isinstance(slide_index, int) and 0 <= slide_index < len(prs.slides):
                rId = prs.slides._sldIdLst[slide_index].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[slide_index]
            else:
                return json.dumps({"error": f"Noto'g'ri slide_index: {slide_index}"})

        prs.save(str(path))
        return json.dumps({
            "status": "ok",
            "message": f"{path.name} yangilandi",
            "action": action,
            "total_slides": len(prs.slides),
        })

    # Register tools
    registry.register(
        name="create_pptx",
        description=(
            "PowerPoint (.pptx) prezentatsiya yaratish. Sarlavha va slaydlar bering. "
            "Har bir slayd title va content dan iborat."
        ),
        parameters={
            "type": "object",
            "required": ["filename"],
            "properties": {
                "filename": {"type": "string", "description": "Fayl nomi (masalan: taqdimot.pptx)"},
                "title": {"type": "string", "description": "Prezentatsiya sarlavhasi"},
                "subtitle": {"type": "string", "description": "Sarlavha ostidagi matn"},
                "slides": {
                    "type": "array",
                    "description": "Slaydlar ro'yxati: [{title, content}]",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "content": {"type": "string"},
                        },
                    },
                },
            },
        },
        handler=create_pptx,
    )

    registry.register(
        name="read_pptx",
        description="PowerPoint (.pptx) prezentatsiyani o'qish. Barcha slaydlar matnini qaytaradi.",
        parameters={
            "type": "object",
            "required": ["file"],
            "properties": {
                "file": {"type": "string", "description": "Fayl nomi yoki to'liq path"},
            },
        },
        handler=read_pptx,
    )

    registry.register(
        name="edit_pptx",
        description=(
            "PowerPoint (.pptx) prezentatsiyani tahrirlash. Slayd qo'shish (add_slide), "
            "matn almashtirish (replace), slayd o'chirish (delete_slide)."
        ),
        parameters={
            "type": "object",
            "required": ["file"],
            "properties": {
                "file": {"type": "string", "description": "Fayl nomi yoki to'liq path"},
                "action": {
                    "type": "string",
                    "enum": ["add_slide", "replace", "delete_slide"],
                    "description": "Harakat turi: add_slide, replace, delete_slide",
                },
                "title": {"type": "string", "description": "Slayd sarlavhasi (add_slide uchun)"},
                "content": {"type": "string", "description": "Slayd matni (add_slide uchun)"},
                "old_text": {"type": "string", "description": "Almashtirilishi kerak bo'lgan matn (replace uchun)"},
                "new_text": {"type": "string", "description": "Yangi matn (replace uchun)"},
                "slide_index": {
                    "type": "integer",
                    "description": "O'chiriladigan slayd indeksi, 0 dan boshlanadi (delete_slide uchun)",
                },
            },
        },
        handler=edit_pptx,
    )

    logger.info("PPTX tools registered: create_pptx, read_pptx, edit_pptx")
